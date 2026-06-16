"""
pptx_lint.py — OBJECTIVE-DEFECT linter for .pptx decks.

SCOPE PHILOSOPHY (read before adding a check):
  This linter only flags OBJECTIVE, UNAMBIGUOUS mistakes — defects that are wrong
  regardless of taste or context, that a strictly MORE capable model would also
  never want to produce. It must never become a stylistic shackle on a better
  model. Every check must pass BOTH tests or it does not belong here:
    (1) OBJECTIVE   — a violation is wrong independent of style/intent.
    (2) NO-SHACKLE  — a smarter model would still always avoid it.
  Therefore the following are DELIBERATELY NOT CHECKED (model capability owns them):
  word/bullet density, font-family count, color palette, layout elegance, margins,
  alignment grids, aspect-ratio choice. Those are "raising the ceiling", not
  "removing an obvious mistake".

Each finding: {id, severity, slide, shape, message}. Severity = ERROR | WARN | INFO.
WARN items marked "(verify on render)" are conservative heuristics — the JUDGMENT
layer (pptx_render.py + an eye) is the ground truth for those.

Exit code: 1 if any ERROR (or any WARN when --strict). 0 otherwise. 2 on usage error.

Usage:
    py pptx_lint.py deck.pptx [--json report.json] [--strict] [--quiet]
"""
import argparse
import json
import math
import os
import re
import sys

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400

# ---- thresholds (one place; only thresholds for OBJECTIVE defects live here) ----
CFG = {
    "offslide_tol_in": 0.05,      # past the edge by more than this = off-canvas
    "overlap_min_frac": 0.15,     # text-on-text intersection > frac of smaller box
    "distort_tol": 0.04,          # |displayed_ratio/native_ratio - 1| above this = stretched
    "unreadable_pt": 8.0,         # explicit font strictly below this = unreadable (not a style floor)
    "invisible_contrast": 2.0,    # text-vs-fill contrast below this = effectively invisible
    "overflow_ratio": 1.20,       # estimated text height > box height * this = likely overflow
    "default_font_pt": 18.0,      # assumed size when none is explicit (estimate only)
}

MD_PATTERNS = [
    (re.compile(r"\*\*[^*]+\*\*"), "literal **bold** markdown"),
    (re.compile(r"__[^_]+__"), "literal __bold__ markdown"),
    (re.compile(r"^\s{0,3}#{1,6}\s"), "literal #-heading markdown"),
    (re.compile(r"^\s{0,3}[-*+]\s+\S"), "literal dash/asterisk bullet"),
    (re.compile(r"\|.+\|.+\|"), "literal pipe-table markdown"),
    (re.compile(r"\[[^\]]+\]\([^)]+\)"), "literal [link](url) markdown"),
]
PLACEHOLDER_TEXT = re.compile(
    r"click to add|lorem ipsum|your text here|insert text here|\[?placeholder\]?|"
    r"sample text|todo:|tbd\b|xxx+",
    re.I,
)


# ---------- geometry helpers ----------
def to_in(emu):
    return None if emu is None else emu / EMU_PER_INCH


def rect_of(shape):
    if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
        return None
    l, t = to_in(shape.left), to_in(shape.top)
    return (l, t, l + to_in(shape.width), t + to_in(shape.height))


def intersect_area(a, b):
    ix = max(0.0, min(a[2], b[2]) - max(a[0], b[0]))
    iy = max(0.0, min(a[3], b[3]) - max(a[1], b[1]))
    return ix * iy


def area(r):
    return max(0.0, r[2] - r[0]) * max(0.0, r[3] - r[1])


def is_cjk(ch):
    o = ord(ch)
    return (0xAC00 <= o <= 0xD7A3 or 0x3000 <= o <= 0x303F or 0x3040 <= o <= 0x30FF
            or 0x3400 <= o <= 0x4DBF or 0x4E00 <= o <= 0x9FFF or 0xFF00 <= o <= 0xFFEF)


# ---------- walkers ----------
def walk_text_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_text_shapes(sh.shapes)
        elif sh.has_text_frame:
            yield sh


def runs_of(shape):
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            yield para, run


def shape_label(shape):
    try:
        return shape.name
    except Exception:
        return "?"


# ---------- color / contrast ----------
def _lin(c):
    c = c / 255.0
    return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4


def luminance(rgb):
    r, g, b = rgb
    return 0.2126 * _lin(r) + 0.7152 * _lin(g) + 0.0722 * _lin(b)


def contrast_ratio(c1, c2):
    l1, l2 = luminance(c1), luminance(c2)
    hi, lo = max(l1, l2), min(l1, l2)
    return (hi + 0.05) / (lo + 0.05)


def explicit_rgb(color):
    try:
        from pptx.enum.dml import MSO_COLOR_TYPE
        if color and color.type == MSO_COLOR_TYPE.RGB:
            rgb = color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None
    return None


def shape_fill_rgb(shape):
    try:
        from pptx.enum.dml import MSO_FILL_TYPE
        fill = shape.fill
        if fill.type == MSO_FILL_TYPE.SOLID:
            return explicit_rgb(fill.fore_color)
    except Exception:
        return None
    return None


# ---------- text-overflow estimate (A2: text spills its own box) ----------
def estimate_overflow(shape):
    """Conservative estimate: does the text need more height than its box has?
    Returns (needed_in, box_in) if likely overflow, else None. Heuristic -> WARN only."""
    if not shape.has_text_frame:
        return None
    tf = shape.text_frame
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        if tf.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
            return None  # PowerPoint shrinks text to fit -> no overflow
    except Exception:
        pass
    box_w = to_in(shape.width)
    box_h = to_in(shape.height)
    if not box_w or not box_h:
        return None
    usable_w = max(0.15, box_w - 0.2)   # ~0.1in internal L/R margins
    usable_h = max(0.1, box_h - 0.1)    # ~0.05in T/B
    word_wrap = tf.word_wrap is not False
    needed_h = 0.0
    for para in tf.paragraphs:
        pt = None
        for run in para.runs:
            if run.font.size is not None:
                pt = run.font.size.pt
                break
        if pt is None and para.font.size is not None:
            pt = para.font.size.pt
        if pt is None:
            pt = CFG["default_font_pt"]
        text = para.text
        if not text.strip():
            needed_h += pt / 72.0 * 1.2
            continue
        width_in = sum((pt / 72.0) * (1.0 if is_cjk(ch) else 0.5) for ch in text)
        lines = 1 if not word_wrap else max(1, math.ceil(width_in / usable_w))
        needed_h += lines * (pt / 72.0 * 1.2)
    if needed_h > usable_h * CFG["overflow_ratio"]:
        return (needed_h, usable_h)
    return None


# ---------- linter ----------
class Linter:
    def __init__(self, cfg=CFG):
        self.cfg = cfg
        self.findings = []

    def add(self, sev, cid, slide, shape, msg):
        self.findings.append({"id": cid, "severity": sev, "slide": slide,
                              "shape": shape, "message": msg})

    def run(self, path):
        prs = Presentation(path)
        W, H = to_in(prs.slide_width), to_in(prs.slide_height)
        for idx, slide in enumerate(prs.slides, start=1):
            shapes = list(slide.shapes)
            self.check_offslide(idx, shapes, W, H)
            self.check_text_on_text(idx, shapes)
            self.check_images(idx, shapes)
            self.check_text(idx, shapes)
        return {
            "file": os.path.abspath(path),
            "slide_size_in": [round(W, 3), round(H, 3)],
            "summary": self.summary(),
            "findings": sorted(self.findings, key=lambda f: (f["slide"] or 0, f["id"])),
        }

    def summary(self):
        s = {"ERROR": 0, "WARN": 0, "INFO": 0}
        for f in self.findings:
            s[f["severity"]] += 1
        return s

    # A1 — content off the canvas
    def check_offslide(self, idx, shapes, W, H):
        tol = self.cfg["offslide_tol_in"]
        for sh in shapes:
            r = rect_of(sh)
            if not r:
                continue
            l, t, rt, b = r
            if l < -tol or t < -tol or rt > W + tol or b > H + tol:
                self.add("ERROR", "GEO-OFFSLIDE", idx, shape_label(sh),
                         f"Off the slide canvas (box [{l:.2f},{t:.2f},{rt:.2f},{b:.2f}] "
                         f"vs slide {W:.2f}x{H:.2f}).")

    # A3 — two text shapes that BOTH carry text overlapping = illegible stack
    def check_text_on_text(self, idx, shapes):
        boxes = []
        for sh in walk_text_shapes(shapes):
            if not sh.text_frame.text.strip():
                continue
            r = rect_of(sh)
            if r and area(r) > 0:
                boxes.append((shape_label(sh), r))
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                ni, ri = boxes[i]
                nj, rj = boxes[j]
                inter = intersect_area(ri, rj)
                if inter <= 0:
                    continue
                frac = inter / min(area(ri), area(rj))
                if frac >= self.cfg["overlap_min_frac"]:
                    self.add("WARN", "GEO-OVERLAP", idx, f"{ni} / {nj}",
                             f"Two text boxes overlap {frac*100:.0f}% of the smaller "
                             f"(text-on-text; verify intentional).")

    # A5 — image stretched off its native aspect ratio
    def check_images(self, idx, shapes):
        for sh in shapes:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            if sh.width is None or sh.height is None or sh.height == 0:
                continue
            try:
                nw, nh = sh.image.size
            except Exception:
                continue
            if not nw or not nh:
                continue
            disp = sh.width / sh.height
            native = nw / nh
            if native == 0:
                continue
            dev = abs(disp / native - 1.0)
            if dev > self.cfg["distort_tol"]:
                self.add("WARN", "IMG-DISTORT", idx, shape_label(sh),
                         f"Image stretched: displayed ratio {disp:.2f} vs native "
                         f"{native:.2f} ({dev*100:.0f}% off). Lock aspect ratio.")

    # A2 + artifacts + narrow legibility checks
    def check_text(self, idx, shapes):
        for sh in walk_text_shapes(shapes):
            # A2 text overflows its own box (estimate)
            ov = estimate_overflow(sh)
            if ov:
                self.add("WARN", "TXT-OVERFLOW", idx, shape_label(sh),
                         f"Text likely overflows its box (needs ~{ov[0]:.1f}in, box "
                         f"~{ov[1]:.1f}in; verify on render).")
            # empty placeholder left in
            if sh.is_placeholder and not sh.text_frame.text.strip():
                self.add("WARN", "ART-PLACEHOLDER", idx, shape_label(sh),
                         "Empty placeholder left on slide.")
            fill_rgb = shape_fill_rgb(sh)
            for para, run in runs_of(sh):
                text = run.text or ""
                # unreadable (truly tiny, explicit only)
                sz = run.font.size
                if sz is not None and sz.pt < self.cfg["unreadable_pt"]:
                    self.add("WARN", "TXT-TINY", idx, shape_label(sh),
                             f"{sz.pt:.0f}pt text is unreadable (< {self.cfg['unreadable_pt']:.0f}pt).")
                # markdown artifact
                for pat, why in MD_PATTERNS:
                    if pat.search(text):
                        self.add("WARN", "ART-MARKDOWN", idx, shape_label(sh),
                                 f"{why}: {text.strip()[:50]!r}")
                        break
                # placeholder / boilerplate phrase
                if PLACEHOLDER_TEXT.search(text):
                    self.add("WARN", "ART-PLACEHOLDER", idx, shape_label(sh),
                             f"Placeholder/boilerplate text: {text.strip()[:50]!r}")
                # near-invisible text (both explicit RGB only)
                run_rgb = explicit_rgb(run.font.color)
                if run_rgb and fill_rgb:
                    cr = contrast_ratio(run_rgb, fill_rgb)
                    if cr < self.cfg["invisible_contrast"]:
                        self.add("WARN", "TXT-INVISIBLE", idx, shape_label(sh),
                                 f"Text nearly invisible on its fill (contrast {cr:.1f}:1).")


# ---------- reporting ----------
SEV_ORDER = {"ERROR": 0, "WARN": 1, "INFO": 2}
SEV_MARK = {"ERROR": "x", "WARN": "!", "INFO": "i"}


def print_human(report):
    s = report["summary"]
    print(f"\nPPTX LINT  {os.path.basename(report['file'])}  "
          f"({report['slide_size_in'][0]}x{report['slide_size_in'][1]} in)")
    print(f"  {s['ERROR']} ERROR   {s['WARN']} WARN   {s['INFO']} INFO\n")
    if not report["findings"]:
        print("  clean -- no objective defects found.\n")
        return
    last = None
    for f in sorted(report["findings"],
                    key=lambda x: (x["slide"] or 0, SEV_ORDER[x["severity"]], x["id"])):
        head = "deck-wide" if f["slide"] is None else f"slide {f['slide']}"
        if head != last:
            print(f"-- {head} --")
            last = head
        loc = f" [{f['shape']}]" if f["shape"] else ""
        print(f"  [{SEV_MARK[f['severity']]}] {f['id']}{loc}: {f['message']}")
    print()


def main():
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
        sys.stderr.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--json", default=None)
    ap.add_argument("--strict", action="store_true", help="exit nonzero on WARN too")
    ap.add_argument("--quiet", action="store_true")
    args = ap.parse_args()

    if not os.path.isfile(args.pptx):
        print("ERROR: no such file:", args.pptx, file=sys.stderr)
        sys.exit(2)

    report = Linter().run(args.pptx)
    if args.json:
        with open(args.json, "w", encoding="utf-8") as fh:
            json.dump(report, fh, ensure_ascii=False, indent=2)
        print("wrote", args.json)
    if not args.quiet:
        print_human(report)

    s = report["summary"]
    sys.exit(1 if s["ERROR"] > 0 or (args.strict and s["WARN"] > 0) else 0)


if __name__ == "__main__":
    main()
