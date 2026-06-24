"""
make_bad_powerpoint.py — build a deliberately flawed .pptx to validate tools/check_powerpoint.py.

Each slide injects KNOWN defects so the linter's deterministic checks can be
confirmed to fire (true positives) without firing on the clean control slide.

Usage:
    python examples/make_bad_powerpoint.py [out.pptx]
Default output: ./_test_flawed_deck.pptx (plus a generated sample image).
"""
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))


def make_sample_image(path, w_px=400, h_px=100, color=(40, 90, 200)):
    """A 4:1 native-ratio image so we can place it squished to 1:1 = distortion."""
    Image.new("RGB", (w_px, h_px), color).save(path)
    return path


def build(out_path):
    prs = Presentation()
    # 16:9 canvas
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    title_only = prs.slide_layouts[5]

    def textbox(slide, left, top, width, height, text, size=18, font="Calibri"):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.name = font
        return tb

    # --- Slide 1: CLEAN control (should pass) ---
    s = prs.slides.add_slide(title_only)
    s.shapes.title.text = "Clean Control Slide"
    textbox(s, Inches(1), Inches(2), Inches(8), Inches(3),
            "One clear idea.\nReadable 18pt body.\nNothing overlaps.", size=18, font="Calibri")

    # --- Slide 2: OFF-SLIDE elements ---
    s = prs.slides.add_slide(blank)
    textbox(s, Inches(11.5), Inches(3), Inches(4), Inches(1),
            "Runs off the right edge", size=18)            # extends to 15.5in > 13.333
    textbox(s, Inches(-2), Inches(1), Inches(3), Inches(1),
            "Starts off the left edge", size=18)            # left negative

    # --- Slide 3: OVERLAPPING text boxes ---
    s = prs.slides.add_slide(blank)
    textbox(s, Inches(2), Inches(2), Inches(6), Inches(2), "Box A occupies this region", size=20)
    textbox(s, Inches(3), Inches(2.4), Inches(6), Inches(2), "Box B sits on top of Box A", size=20)

    # --- Slide 4: TINY font ---
    s = prs.slides.add_slide(blank)
    textbox(s, Inches(1), Inches(3), Inches(10), Inches(1),
            "This 6pt line is unreadable from any seat in the room.", size=6)

    # --- Slide 5: MARKDOWN leakage + empty placeholder ---
    s = prs.slides.add_slide(title_only)
    s.shapes.title.text = "**Key Point**"                   # literal markdown in title
    textbox(s, Inches(1), Inches(2), Inches(10), Inches(3),
            "## Heading left as text\n- bullet written with a dash\n**bold** never rendered\n| col | col |",
            size=18)
    # leave a body placeholder empty on a layout that has one (title_only has none extra),
    # so add an explicit empty textbox to simulate "Click to add text"
    empty = s.shapes.add_textbox(Inches(1), Inches(5.5), Inches(4), Inches(0.8))
    empty.text_frame.text = "Click to add text"

    # --- Slide 6: IMAGE distortion (4:1 native placed as 1:1) ---
    s = prs.slides.add_slide(blank)
    img = make_sample_image(os.path.join(HERE, "_sample_4x1.png"))
    s.shapes.add_picture(img, Inches(4), Inches(2), Inches(3), Inches(3))  # squished to square

    # --- Slide 7: TEXT OVERFLOWS ITS BOX (small box, big text) [A2] ---
    s = prs.slides.add_slide(blank)
    long_para = ("This 24pt paragraph is far too long for the tiny box it was placed "
                 "in, so it spills past the bottom edge of its own text frame and gets "
                 "clipped on screen instead of fitting cleanly inside the shape.")
    textbox(s, Inches(1), Inches(3), Inches(3), Inches(0.8), long_para, size=24)

    # --- Slide 8: NEAR-INVISIBLE text + multiple fonts (fonts must NOT be flagged) ---
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(2), Inches(3), Inches(8), Inches(1))
    tb.fill.solid()
    tb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)      # white fill
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "Nearly invisible: light gray text on a white fill"
    r.font.size = Pt(20)
    r.font.color.rgb = RGBColor(0xF2, 0xF2, 0xF2)            # ~white text -> ~1.05:1
    # three different fonts on one slide -> CONTROL: must produce NO finding
    textbox(s, Inches(2), Inches(4.4), Inches(8), Inches(1), "Calibri / Times / Courier mix is fine",
            size=22, font="Times New Roman")

    # --- Slide 9: DENSE but FITTING content [CONTROL] -> must produce NO finding ---
    # proves the linter does not shackle information density (model capability owns that)
    s = prs.slides.add_slide(blank)
    dense = "\n".join([f"Point {i}: a substantive, deliberately information-rich bullet "
                       f"that a capable deck may legitimately include." for i in range(1, 11)])
    textbox(s, Inches(0.6), Inches(0.6), Inches(12.1), Inches(6.3), dense, size=14)

    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else os.path.join(HERE, "_test_flawed_deck.pptx")
    p = build(out)
    print("wrote", p)
